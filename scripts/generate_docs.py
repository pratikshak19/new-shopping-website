#!/usr/bin/env python3
"""Generate college-ready PDF, Word and PowerPoint files for Trendora."""
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "public" / "docs"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = (26, 26, 46)
ROSE = (255, 63, 108)

CHAPTERS = [
    (
        "Certificate",
        "This is to certify that the project titled “Trendora — A Multi-Authority Online Shopping Website” "
        "is a bona fide record of work carried out as part of the academic curriculum. The software implements "
        "a customer storefront inspired by Myntra, Meesho and Flipkart together with Seller, Admin and Owner consoles. "
        "The candidate has presented the working system, test cases and this report for evaluation.",
    ),
    (
        "Declaration",
        "I declare that this project is my original work. Public storefronts of Myntra, Flipkart and Meesho were "
        "studied only as user-experience references. Brand name, source code, catalogue copy and documents are original. "
        "No live payment gateway is connected. Simulated credentials are published for laboratory demonstration only.",
    ),
    (
        "Acknowledgement",
        "I thank my project guide, the Department of Computer Science / Information Technology, laboratory staff "
        "and classmates who reviewed the four-role login flow. I also thank open documentation of React, Vite and React Router.",
    ),
    (
        "Abstract",
        "Indian e-commerce is defined by three consumer habits: fashion discovery (Myntra), price and deals (Flipkart) "
        "and supplier-led catalogue (Meesho). Trendora unifies those habits in one React single-page application and "
        "adds four login authorities — Customer, Seller, Admin, Owner — each with a separate home and permission set.\n\n"
        "Customers browse eight categories, search and filter, read size charts, check PIN-code serviceability, "
        "compare up to three SKUs, heart a wishlist, apply coupons, gift-wrap, check out (UPI / card / COD simulated), "
        "track, cancel and return orders, earn Insider points, and open help tickets. Sellers list SKUs and see "
        "commissioned earnings. Admins run catalogue, order pipeline, returns, coupons and the help desk. "
        "The Owner additionally controls people, roles, store policy and GMV reports.\n\n"
        "State lives in a Context provider persisted to localStorage (key trendora-store-v2) so a viva demo survives refresh "
        "without a paid back-end. This report documents objectives, literature, SRS, design, implementation, testing and future scope.",
    ),
    (
        "1. Introduction",
        "1.1 Motivation\nCollege students already shop on Meesho, Myntra and Flipkart. Cloning only the storefront leaves "
        "the authority model unexplained. Real platforms separate the shopper, the seller, the operations admin and the owner. "
        "Trendora therefore implements all four logins.\n\n"
        "1.2 Problem statement\nDesign a responsive multi-role shopping website that covers the shopping journey of Indian "
        "fashion/general e-commerce and the back-office of a marketplace, without a live warehouse or payment gateway.\n\n"
        "1.3 Objectives\n"
        "1. Study IA of Myntra, Meesho and Flipkart.\n"
        "2. Implement customer features: catalogue, search, filter, sort, PDP, size chart, PIN check, reviews, Q&A, "
        "wishlist, compare, bag, coupons, gift wrap, checkout, orders, cancel, return, addresses, Insider, studio, offers, help.\n"
        "3. Implement seller studio: list SKUs, own orders, earnings after commission.\n"
        "4. Implement admin desk: catalogue, order pipeline, returns, coupons, tickets, user block.\n"
        "5. Implement owner console: all admin rights plus create staff, change roles, store settings, GMV reports.\n"
        "6. Persist session data locally and ship PDF / Word / PPT documents.\n\n"
        "1.4 Scope\nIn scope: the modules above, INR pricing, simulated payments, college documents.\n"
        "Out of scope: real Razorpay capture, AWB logistics, native apps, ML recommendations.",
    ),
    (
        "2. Literature survey",
        "Myntra — fashion PDP, size chart, Insider loyalty, Studio lookbooks, easy returns.\n"
        "Flipkart — deal timer, coupon engine, MRP break-up, COD, order timeline.\n"
        "Meesho — supplier listing, commission payout, category-first mobile bag.\n"
        "Amazon Seller Central / Flipkart Seller — inventory, order states, returns desk.\n\n"
        "Gap: a student cannot reimplement logistics. Trendora copies the nouns (bag, wishlist, COD, commission) "
        "and implements them as original React code with an honest simulation note.",
    ),
    (
        "3. System analysis",
        "3.1 Actors and authority\n"
        "Customer — shop, bag, pay, track, return, review, ticket.\n"
        "Seller — list own products, view orders containing those SKUs, see net payout.\n"
        "Admin — all catalogue, all orders, returns, coupons, tickets, block users. Cannot change owner or store policy.\n"
        "Owner — superset of admin + create staff + change roles + settings (commission, free-ship, announcement) + reports.\n\n"
        "3.2 Functional requirements (selected)\n"
        "FR1 Multi-category catalogue with sellerId.\nFR2 Search / filter / sort.\nFR3 PDP variants, size chart, PIN.\n"
        "FR4 Bag, wishlist, compare (max 3).\nFR5 Coupon engine with min-cart and Insider gate.\n"
        "FR6 Checkout with saved addresses, gift wrap, UPI/card/COD.\nFR7 Order timeline + cancel + return.\n"
        "FR8 Reviews and Q&A.\nFR9 Four-role auth, block, role change.\nFR10 Seller CRUD + earnings.\n"
        "FR11 Admin/Owner operations desks.\nFR12 Help tickets.\nFR13 Notifications and Insider points.\n\n"
        "3.3 Non-functional\nResponsive, INR format, no crash on empty bag, honest security disclaimer, lab-machine friendly.",
    ),
    (
        "4. System design",
        "4.1 Architecture\nBrowser → React pages → StoreContext (use cases) → products.js + localStorage.\n"
        "Three layers: presentation, application state, persistence.\n\n"
        "4.2 Routing\nPublic: /, /shop, /product/:id, /offers, /brands, /studio, /insider, /compare, /help, /login, /register.\n"
        "Customer: /cart, /checkout, /orders, /returns, /addresses, /wishlist, /profile, /notifications.\n"
        "Guarded: /seller/*, /admin/*, /owner/* via RequireRole.\n\n"
        "4.3 Price algorithm\n"
        "grand = max(0, subtotal − coupon + shipping + giftWrap)\n"
        "shipping = 0 if subtotal >= freeShipMin else shipFee (owner configurable; default 999 / 79).\n"
        "giftWrap = 49 if selected.\n"
        "Insider points += floor(grand / 10).\n"
        "Seller net = item sale − commission% (default 12).\n\n"
        "4.4 Order state machine\nConfirmed → Packed → Shipped → Out for delivery → Delivered → (optional) Returned.\n"
        "Cancel allowed before Shipped. Return allowed only after Delivered.\n\n"
        "4.5 Demo accounts\n"
        "demo@trendora.in / demo123 (Customer)\n"
        "seller@trendora.in / seller123 (Seller)\n"
        "admin@trendora.in / admin123 (Admin)\n"
        "owner@trendora.in / owner123 (Owner)",
    ),
    (
        "5. Implementation",
        "Environment: Node.js 18+, npm, Vite 5, React 18, React Router 6, Context API, CSS design tokens.\n\n"
        "Key modules:\n"
        "src/context/StoreContext.jsx — every mutation (auth, bag, orders, catalogue, tickets).\n"
        "src/data/products.js — seed catalogue, coupons, size charts, PIN helper.\n"
        "src/pages/dash/* — Owner / Admin / Seller consoles.\n"
        "src/components/RequireRole.jsx — route guard.\n\n"
        "Security honesty: demo passwords sit in localStorage. Production would hash on a server (bcrypt) and issue JWT/cookies. "
        "This is stated in FAQ, slides and viva notes.",
    ),
    (
        "6. Testing",
        "T1 Home hero + categories render — Pass\n"
        "T2 Search “earbuds” finds PulseBuds — Pass\n"
        "T3 Filter Women + sort price — Pass\n"
        "T4 Two sizes of one SKU become two bag lines — Pass\n"
        "T5 FESTIVE20 rejected below Rs 1999 — Pass\n"
        "T6 PIN 56A rejected, 416001 accepted — Pass\n"
        "T7 Customer places order; bag empties; stock drops — Pass\n"
        "T8 Admin marks Packed → Shipped → Delivered — Pass\n"
        "T9 Customer requests return; admin refunds — Pass\n"
        "T10 Seller lists a SKU; it appears on /shop — Pass\n"
        "T11 Owner creates an admin; new login works — Pass\n"
        "T12 Blocked user cannot sign in — Pass\n"
        "T13 Refresh keeps session (trendora-store-v2) — Pass\n"
        "T14 /owner as customer redirects home — Pass",
    ),
    (
        "7. Results and conclusion",
        "Trendora presents a credible Indian storefront and three staff consoles. A full purchase plus an admin status change "
        "takes under three minutes in a viva. Context + localStorage is a pedagogical stand-in for REST; swapping placeOrder "
        "for fetch('/api/orders') is a weekend of work, not a rewrite.\n\n"
        "Conclusion: the project meets its objectives. It is demoable, documented (PDF, Word, PPT) and honest about simulation.",
    ),
    (
        "8. Future scope",
        "1. Express + MongoDB API.\n2. Razorpay / Stripe test mode.\n3. Image upload for sellers.\n"
        "4. Size recommender.\n5. PWA offline catalogue.\n6. Jest + React Testing Library.\n7. Real SMS/email OTPs.",
    ),
    (
        "9. References",
        "1. React documentation — https://react.dev\n"
        "2. Vite guide — https://vitejs.dev\n"
        "3. React Router — https://reactrouter.com\n"
        "4. Nielsen Norman Group, e-commerce UX heuristics\n"
        "5. Public storefronts of Myntra, Flipkart and Meesho (UX reference only)\n"
        "6. MDN Web Docs — Web Storage API",
    ),
    (
        "Appendix A — How to run",
        "cd new-shopping-website\nnpm install\nnpm run dev\nOpen the printed URL (default http://localhost:5173).\n"
        "Production: npm run build && npm run preview.",
    ),
    (
        "Appendix B — User manuals by role",
        "Customer: Login demo@trendora.in → shop → PDP → bag → FESTIVE20 → checkout → track.\n"
        "Seller: Login seller@trendora.in → Catalogue → list SKU → Earnings.\n"
        "Admin: Login admin@trendora.in → Orders → mark Packed / Shipped / Delivered → Returns / Tickets.\n"
        "Owner: Login owner@trendora.in → People & roles → create staff → Settings → Reports.",
    ),
]

SLIDES = [
    ("Trendora", "A multi-authority shopping website\nMeesho · Myntra · Flipkart\nCollege / vlg project 2025–26"),
    ("Why this project", "Everyone already shops online.\nE-commerce teaches data + money + identity + UX.\nFour logins make the viva complete."),
    ("Problem", "Build discover → decide → pay → track\nand a seller / admin / owner back-office\nwithout a real warehouse."),
    ("Inspiration", "Myntra — PDP, size chart, Insider, Studio, returns\nFlipkart — deals, coupons, COD, timeline\nMeesho — seller listing + commission"),
    ("Four authorities", "Customer — shop\nSeller — catalogue + earnings\nAdmin — operations desk\nOwner — people, policy, GMV"),
    ("Customer features", "Search, filter, sort, size chart, PIN check\nReviews, Q&A, wishlist, compare\nBag, coupons, gift wrap, checkout\nOrders, cancel, return, Insider, help"),
    ("Seller studio", "List / edit / hide SKUs\nSee only own order lines\nPayout = sale − commission"),
    ("Admin desk", "All catalogue\nOrder pipeline\nReturns & refunds\nCoupons & tickets\nBlock users"),
    ("Owner console", "Everything admin can do\nCreate admin / seller logins\nChange roles\nCommission, free-ship, announcement\nReports"),
    ("Architecture", "UI pages\n   ↓\nStoreContext (use cases)\n   ↓\nproducts.js + localStorage v2"),
    ("Money formula", "grand = subtotal − coupon + ship + wrap\nFree ship ≥ Rs 999 (owner can change)\n1 Insider point / Rs 10"),
    ("Order machine", "Confirmed → Packed → Shipped\n→ Out for delivery → Delivered\nCancel before ship · Return after deliver"),
    ("Live demo script", "1. Customer buys a dress + FESTIVE20\n2. Admin marks Delivered\n3. Customer requests return\n4. Seller adds a SKU\n5. Owner creates a staff login"),
    ("Testing", "14 cases: search, coupon min, PIN,\nstock drop, pipeline, return,\nrole guard, block, refresh."),
    ("Be honest", "No real money.\nPasswords not hashed on a server.\nNothing ships.\nSimulation with production UX."),
    ("Documents", "Project report — PDF + Word\nHow it was built — PDF\nPresentation — PPTX + HTML"),
    ("Thank you", "Questions?\nStore · /login role cards · /documents"),
]


class ReportPDF(FPDF):
    def header(self):
        if self.page_no() == 1:
            return
        self.set_font("Helvetica", "I", 9)
        self.set_text_color(*ROSE)
        self.cell(0, 8, "Trendora  |  College project report", ln=1)
        self.set_draw_color(*NAVY)
        self.line(15, 16, 195, 16)
        self.ln(4)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 8, f"Page {self.page_no()}", align="C")


def ascii(s: str) -> str:
    table = str.maketrans({
        "\u2014": "-",
        "\u2013": "-",
        "\u2212": "-",
        "\u201c": '"',
        "\u201d": '"',
        "\u2018": "'",
        "\u2019": "'",
        "\u2022": "-",
        "\u2192": "-",
        "\u2193": "-",
        "\u00d7": "x",
        "\u20b9": "Rs ",
        "\u00b7": "-",
        "\u2265": ">=",
    })
    return s.translate(table).encode("latin-1", "replace").decode("latin-1")


def write_pdf(path: Path, title: str, subtitle: str, chapters):
    pdf = ReportPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()
    pdf.set_fill_color(*NAVY)
    pdf.rect(0, 0, 210, 297, "F")
    pdf.set_text_color(*ROSE)
    pdf.set_font("Helvetica", "B", 13)
    pdf.set_xy(20, 70)
    pdf.cell(0, 10, "COLLEGE / VLG PROJECT  2025-26")
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "B", 32)
    pdf.set_xy(20, 90)
    pdf.multi_cell(170, 14, ascii(title))
    pdf.set_font("Helvetica", "", 14)
    pdf.set_xy(20, 140)
    pdf.multi_cell(170, 8, ascii(subtitle))
    pdf.set_font("Helvetica", "", 11)
    pdf.set_xy(20, 250)
    pdf.cell(0, 8, "Department of Computer Science / Information Technology")

    pdf.add_page()
    pdf.set_text_color(*NAVY)
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "Contents", ln=1)
    pdf.set_font("Helvetica", "", 11)
    for i, (h, _) in enumerate(chapters, 1):
        pdf.cell(0, 7, ascii(f"{i}.  {h}"), new_x="LMARGIN", new_y="NEXT")

    for heading, body in chapters:
        pdf.add_page()
        pdf.set_text_color(*ROSE)
        pdf.set_font("Helvetica", "B", 16)
        pdf.multi_cell(0, 9, ascii(heading))
        pdf.ln(2)
        pdf.set_text_color(40, 40, 55)
        pdf.set_font("Helvetica", "", 11)
        for para in body.split("\n"):
            pdf.multi_cell(0, 6.2, ascii(para) if para else " ")
            pdf.ln(1)
    pdf.output(path)


def write_docx(path: Path):
    doc = Document()
    styles = doc.styles["Normal"]
    styles.font.name = "Calibri"
    styles.font.size = Pt(11)

    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run("TRENDORA")
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(*NAVY)

    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = s.add_run("A Multi-Authority Online Shopping Website\nInspired by Meesho, Myntra and Flipkart")
    r.font.size = Pt(14)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(
        "Project Report\nSubmitted in partial fulfilment of the requirements\n"
        "of the college / vlg programme\nAcademic year 2025–26"
    )

    doc.add_page_break()
    doc.add_heading("Table of contents", level=1)
    for i, (h, _) in enumerate(CHAPTERS, 1):
        doc.add_paragraph(f"{i}. {h}")

    for heading, body in CHAPTERS:
        doc.add_heading(heading, level=1)
        for para in body.split("\n"):
            doc.add_paragraph(para)

    doc.add_heading("Appendix C — Module list", level=1)
    for item in [
        "Home, Shop, Product detail, Offers, Brands, Studio, Insider, Compare",
        "Bag, Checkout, Orders, Returns, Addresses, Wishlist, Profile, Notifications, Help",
        "Login (4 role cards), Register (customer / seller)",
        "Seller: overview, catalogue, orders, earnings, tickets",
        "Admin: overview, catalogue, orders, returns, users, coupons, tickets",
        "Owner: all admin + reports + settings + create staff",
        "Documents page with PDF / Word / PPT downloads",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.save(path)


def write_pptx(path: Path):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for title, body in SLIDES:
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = PRGB(*NAVY)
        box = slide.shapes.add_textbox(PInches(0.8), PInches(1.6), PInches(11.6), PInches(1.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(40)
        p.font.bold = True
        p.font.color.rgb = PRGB(255, 255, 255)
        p.font.name = "Calibri"
        box2 = slide.shapes.add_textbox(PInches(0.8), PInches(3.2), PInches(11.6), PInches(3.4))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        first = True
        for line in body.split("\n"):
            para = tf2.paragraphs[0] if first else tf2.add_paragraph()
            first = False
            para.text = line
            para.font.size = PPt(22)
            para.font.color.rgb = PRGB(230, 214, 208)
            para.space_after = PPt(8)
        tag = slide.shapes.add_textbox(PInches(0.8), PInches(0.4), PInches(6), PInches(0.4))
        tp = tag.text_frame.paragraphs[0]
        tp.text = "TRENDORA  ·  COLLEGE PROJECT"
        tp.font.size = PPt(12)
        tp.font.color.rgb = PRGB(*ROSE)
    prs.save(path)


def main():
    write_pdf(
        OUT / "Trendora_Project_Report.pdf",
        "Trendora",
        "A multi-authority online shopping website\ninspired by Meesho, Myntra and Flipkart\n\nProject report  ·  PDF copy for submission",
        CHAPTERS,
    )
    write_pdf(
        OUT / "Trendora_How_It_Was_Built.pdf",
        "How Trendora was built",
        "Step-by-step construction log for the viva.\nReact 18 · Vite 5 · Context API · four login authorities.",
        [
            CHAPTERS[0],
            CHAPTERS[4],
            CHAPTERS[7],
            CHAPTERS[8],
            CHAPTERS[13],
            CHAPTERS[14],
        ],
    )
    write_docx(OUT / "Trendora_Project_Report.docx")
    write_pptx(OUT / "Trendora_Presentation.pptx")
    print("Wrote", list(OUT.glob("Trendora_*")))


if __name__ == "__main__":
    main()
